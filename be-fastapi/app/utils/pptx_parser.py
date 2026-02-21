import asyncio
from typing import Any

from app.exceptions import FileParseException


async def parse_pptx(file_path: str) -> tuple[str, dict[str, Any]]:
    try:
        text, metadata = await asyncio.to_thread(_extract_pptx_text, file_path)
        return text, metadata
    except Exception as exc:
        raise FileParseException(f"PPTX parsing failed: {exc}") from exc


def _extract_pptx_text(file_path: str) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides: list[str] = []
    slide_count = 0

    for slide in prs.slides:
        slide_count += 1
        slide_texts: list[str] = []

        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame:
                for para in text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            else:
                table = getattr(shape, "table", None)
                if table:
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))

        if slide_texts:
            slides.append(f"[Slide {slide_count}]\n" + "\n".join(slide_texts))

    return "\n\n".join(slides), {"slide_count": slide_count}
